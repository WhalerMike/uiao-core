"""PPTX leadership briefing generator.

ADR-0005: Produces a styled PowerPoint deck from UIAO canon data,
embedding Mermaid-rendered PNGs, Gemini-generated images, and
matplotlib charts. Designed for agency leadership presentations.
"""
from __future__ import annotations

import logging
from datetime import datetime
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from uiao_core.utils.context import get_settings, load_context

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------
_NAVY = RGBColor(0x1B, 0x3A, 0x5C)
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_ACCENT = RGBColor(0x44, 0x72, 0xC4)
_GRAY = RGBColor(0x80, 0x80, 0x80)
_SLIDE_WIDTH = Inches(13.333)
_SLIDE_HEIGHT = Inches(7.5)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def _add_title_slide(
    prs: Presentation,
    title: str,
    subtitle: str,
) -> None:
    """Add a title slide with navy background."""
    layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(layout)

    # Title
    tf = slide.shapes.title.text_frame
    tf.text = title
    for para in tf.paragraphs:
        para.font.size = Pt(36)
        para.font.color.rgb = _NAVY
        para.font.bold = True
        para.alignment = PP_ALIGN.CENTER

    # Subtitle
    if slide.placeholders[1]:
        stf = slide.placeholders[1].text_frame
        stf.text = subtitle
        for para in stf.paragraphs:
            para.font.size = Pt(18)
            para.font.color.rgb = _ACCENT
            para.alignment = PP_ALIGN.CENTER


def _add_content_slide(
    prs: Presentation,
    title: str,
    bullets: list[str],
) -> Any:
    """Add a content slide with title and bullet points."""
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)

    # Title
    tf = slide.shapes.title.text_frame
    tf.text = title
    for para in tf.paragraphs:
        para.font.size = Pt(28)
        para.font.color.rgb = _NAVY
        para.font.bold = True

    # Body
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, bullet in enumerate(bullets):
        if i == 0:
            body.paragraphs[0].text = bullet
            para = body.paragraphs[0]
        else:
            para = body.add_paragraph()
            para.text = bullet
        para.font.size = Pt(16)
        para.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        para.space_after = Pt(8)

    return slide


def _add_image_slide(
    prs: Presentation,
    title: str,
    image_path: Path,
) -> bool:
    """Add a slide with a centered image."""
    if not image_path.exists():
        logger.warning("Image not found: %s", image_path)
        return False

    layout = prs.slide_layouts[5]  # Blank
    slide = prs.slides.add_slide(layout)

    # Title text box
    from pptx.util import Inches as In

    txBox = slide.shapes.add_textbox(In(0.5), In(0.3), In(12), In(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.color.rgb = _NAVY
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Image centered
    slide.shapes.add_picture(
        str(image_path),
        left=In(1.0),
        top=In(1.3),
        width=In(11.0),
    )
    return True


def _add_table_slide(
    prs: Presentation,
    title: str,
    headers: list[str],
    rows: list[list[str]],
) -> None:
    """Add a slide with a data table."""
    layout = prs.slide_layouts[5]  # Blank
    slide = prs.slides.add_slide(layout)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.color.rgb = _NAVY
    p.font.bold = True

    # Table
    n_rows = min(len(rows) + 1, 15)  # cap at 15 rows per slide
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(0.5), Inches(1.2),
        Inches(12.3), Inches(5.5),
    )
    table = table_shape.table

    # Header row
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for para in cell.text_frame.paragraphs:
            para.font.size = Pt(11)
            para.font.bold = True
            para.font.color.rgb = _WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = _NAVY

    # Data rows
    for r_idx, row_data in enumerate(rows[:n_rows - 1]):
        for c_idx, val in enumerate(row_data):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = str(val)
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(9)


# ---------------------------------------------------------------------------
# Builder
# ---------------------------------------------------------------------------
def _build_pptx(context: dict[str, Any], settings: Any) -> Presentation:
    """Build the leadership briefing PPTX."""
    prs = Presentation()
    prs.slide_width = _SLIDE_WIDTH
    prs.slide_height = _SLIDE_HEIGHT

    lb = context.get("leadership_briefing", {})
    if not isinstance(lb, dict):
        lb = {}

    version = context.get("version", "1.0")
    today = datetime.now().strftime("%B %d, %Y")

    # --- Title Slide ---
    _add_title_slide(
        prs,
        "Unified Identity-Addressing-Overlay\nArchitecture (UIAO)",
        f"Leadership Briefing v{version} | {today} | FedRAMP Moderate",
    )

    # --- Executive Summary ---
    exec_summary = str(lb.get("executive_summary", ""))[:500]
    _add_content_slide(prs, "Executive Summary", [
        s.strip() for s in exec_summary.split(".") if s.strip()
    ][:6])

    # --- Program Vision ---
    vision = str(lb.get("program_vision", ""))[:400]
    _add_content_slide(prs, "Program Vision", [
        s.strip() for s in vision.split(".") if s.strip()
    ][:5])

    # --- Control Planes ---
    planes = lb.get("control_planes", [])
    if planes:
        plane_bullets = []
        for plane in planes:
            if isinstance(plane, dict):
                name = plane.get("name", "")
                desc = str(plane.get("narrative", ""))[:80]
                plane_bullets.append(f"{name}: {desc}")
        _add_content_slide(prs, "The Five Control Planes", plane_bullets[:5])

    # --- Mermaid Diagrams ---
    mermaid_dir = settings.project_root / "assets" / "images" / "mermaid"
    if mermaid_dir.exists():
        for png in sorted(mermaid_dir.glob("*.png")):
            _add_image_slide(prs, png.stem.replace("-", " ").title(), png)

    # --- Gemini AI Visuals ---
    gemini_dir = settings.project_root / "assets" / "images" / "gemini"
    if gemini_dir.exists():
        for png in sorted(gemini_dir.glob("*.png")):
            _add_image_slide(prs, png.stem.replace("-", " ").title(), png)

    # --- Existing Vibrant Visuals ---
    visuals_dir = settings.visuals_dir
    vibrant_images = [
        ("Modernization Atlas", "uiao-vibrant-modernization-atlas.png"),
        ("FedRAMP 20x Governance", "uiao-vibrant-20x-governance-hub.png"),
        ("Identity-to-IP Mapping", "uiao-vibrant-u-plus-a-mapping.png"),
    ]
    for slide_title, img_name in vibrant_images:
        img_path = visuals_dir / img_name
        _add_image_slide(prs, slide_title, img_path)

    # --- Charts ---
    charts_dir = settings.project_root / "assets" / "images"
    radar = charts_dir / "dynamic-maturity-radar.png"
    if radar.exists():
        _add_image_slide(prs, "CISA Zero Trust Maturity Assessment", radar)

    # --- Compliance Matrix Table ---
    matrix = context.get("unified_compliance_matrix", [])
    if isinstance(matrix, list) and matrix:
        headers = ["Pillar", "CISA Pillar", "Maturity", "NIST Controls", "Impact"]
        rows = []
        for entry in matrix[:12]:
            if isinstance(entry, dict):
                controls = entry.get("nist_controls", [])
                ctrl_str = ", ".join(controls) if isinstance(controls, list) else str(controls)
                rows.append([
                    entry.get("pillar", ""),
                    entry.get("cisa_pillar", ""),
                    entry.get("cisa_maturity", ""),
                    ctrl_str,
                    entry.get("impact_statement", "")[:60],
                ])
        if rows:
            _add_table_slide(prs, "Unified Compliance Matrix", headers, rows)

    # --- Closing Slide ---
    _add_title_slide(
        prs,
        "Questions?",
        f"UIAO Program | {today} | CUI",
    )

    return prs


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------
def build_pptx(
    *,
    canon_path: Path | None = None,
    data_dir: Path | None = None,
    exports_dir: Path | None = None,
) -> Path:
    """Generate a leadership briefing PPTX and return the output path."""
    settings = get_settings()
    canon_path = canon_path or settings.canon_dir / "uiao_leadership_briefing_v1.0.yaml"
    data_dir = data_dir or settings.data_dir
    exports_dir = exports_dir or settings.exports_dir

    logger.info("Loading UIAO context for PPTX...")
    context = load_context(canon_path=canon_path, data_dir=data_dir)

    pptx_dir = exports_dir / "pptx"
    pptx_dir.mkdir(parents=True, exist_ok=True)
    out_path = pptx_dir / "UIAO_Leadership_Briefing_v1.0.pptx"

    prs = _build_pptx(context, settings)
    prs.save(str(out_path))

    logger.info("PPTX exported -> %s", out_path)
    return out_path
