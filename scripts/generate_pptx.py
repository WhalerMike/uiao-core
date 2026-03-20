import yaml
import os
from pptx import Presentation
from pptx.util import Inches, Pt

# Paths
CANON_PATH = 'canon/uiao_pitch_deck_v1.0.yaml'
OUTPUT_PATH = 'docs/modernization_atlas_pitch.pptx'
VISUALS_DIR = 'visuals/'


def create_pptx():
    # Load the YAML canon
    with open(CANON_PATH, 'r') as file:
        data = yaml.safe_load(file)

    prs = Presentation()

    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = data['pitch_deck']['title']
    subtitle.text = data['pitch_deck']['subtitle']

    # Loop through slides defined in YAML
    for slide_data in data['pitch_deck']['slides']:
        bullet_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_layout)

        # Set Title
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = slide_data['title']

        # Set Body/Bullets
        body_shape = shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.text = slide_data.get('content', 'Key Objectives:')

        for point in slide_data['points']:
            p = tf.add_paragraph()
            p.text = point
            p.level = 1

        # Placeholder for Visuals
        if 'visual' in slide_data:
            img_path = os.path.join(VISUALS_DIR, slide_data['visual'])
            if os.path.exists(img_path):
                # Add the actual image if it exists
                slide.shapes.add_picture(img_path, Inches(6), Inches(1.5), height=Inches(4))
            else:
                # Add a text box placeholder if image is missing
                txBox = slide.shapes.add_textbox(Inches(6), Inches(2), Inches(3), Inches(1))
                tf_box = txBox.text_frame
                tf_box.text = f"[PLACEHOLDER: Missing Visual {slide_data['visual']}]"

    # Executive Summary Slide
    summary_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(summary_layout)
    slide.shapes.title.text = "Executive Summary"
    tf = slide.placeholders[1].text_frame
    tf.text = data['pitch_deck']['summary']['text']

    prs.save(OUTPUT_PATH)
    print(f"PowerPoint generated successfully at: {OUTPUT_PATH}")


if __name__ == "__main__":
    create_pptx()
